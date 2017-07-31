import os

from datamodels import Lyrics, Song

from hanziconv import HanziConv

from pptx import Presentation

hzc = HanziConv()
# Keep a list of song keys
song_datamodel = list(Song().to_dict().keys())

ALLOWED_EXTENSIONS = set(['pdf'])


def get_lyrics(request, exclude_id=None):
    """
    Utility function that returns a Lyrics object containing the song lyrics.

    Parameters:
    ===========
    - data: (dict-like) pass in `request.form` from Flask app.
    - request: pass in `request` from Flask app.
    - exclude_id: (int) used in excluding a particular lyrics section.
    """
    # Defensive programming checks
    if exclude_id:
        assert isinstance(exclude_id, int)

    # Get lyrics
    l = Lyrics()
    for k, v in request.form.items():
        if 'section' in k:
            idx = int(k.split('-')[-1])
            if idx is not exclude_id:
                lyrics = hzc.toTraditional(request.form[f'lyrics-{idx}'])
                section = request.form[k]
                l.add_section(section=section,
                              lyrics=lyrics)
    return l


def update_song_info(request, eid, song_db, exclude_id=None):
    """
    Updates song information in database.

    Parameters:
    ===========
    - request: (dict-like) the `request` object from the Flask app.
    - eid: (int) the eid of the song to be updated in the database.
    """
    data = {k: hzc.toTraditional(v)
            for k, v in request.form.items()
            if k in song_datamodel}

    lyrics = get_lyrics(request=request, exclude_id=exclude_id)
    data['lyrics'] = lyrics.to_dict()
    song_db.update(data, eids=[eid])


def clean_song_arrangement(arrangement, song_data):
    """
    Cleans the song arrangement and turns it into a list.

    Parameters:
    ===========
    - arrangement: (str) a string containing the arrangement of the song.
    - song_data: (dict) a data dictionary. Keys are the data model fields as
                 specified in `datamodels.py`. One of the keys has to be
                 `lyrics`.

    Returns:
    ========
    - arrangement: (list of str) a list of strings, each of which is a key in
                   song's lyrics dictionary.
    """
    print(arrangement)
    arrangement = arrangement.split(',')
    arrangement = [a.strip(' ') for a in arrangement]
    for section in arrangement:
        assert section in song_data['lyrics'].keys(), (
            f'{section} not present in lyrics.')

    return arrangement


def arrange_lyrics(arrangement, song_data):
    """
    Returns the lyrics of a song arranged by the song data.

    Parameters:
    ===========
    - arrangement: (list of str) a list of strings, each of which is a key in
                   the song's lyrics dictionary.
    - song_data: (dict) the song's data dictionary conforming to the
                 specification in `datamodels.py`. One of the keys is `lyrics`.

    Returns:
    ========
    - arranged_lyrics: (str) the lyrics arranged according to the specified
                       arrangement.
    """
    # Now, we allow the default arrangement to be set.
    print(arrangement)
    arranged_lyrics = ''
    for a in arrangement:
        arranged_lyrics += song_data['lyrics'][a]
        arranged_lyrics += '\n\n'
    print(arranged_lyrics)
    return arranged_lyrics


def make_lyrics_presentation(song_data):
    """
    Makes a set of slides from the lyrics.

    Parameters:
    ===========
    - song_data: (dict) the song's data dictionary conforming to the
                 specification in `datamodels.py`. One of the keys is `lyrics`.
    """
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]
    for sld, lyr in song_data['lyrics'].items():
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = sld
        tf = body_shape.text_frame
        tf.text = lyr
    prs.save('tmp/slides.pptx')


def ensure_dir(file_path):
    directory = os.path.dirname(file_path)
    if not os.path.exists(directory):
        os.makedirs(directory)


def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS
